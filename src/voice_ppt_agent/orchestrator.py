"""Coordinate planning, editing, validation, and retry behavior."""
import json
import os
import subprocess
import sys
import tempfile
from pathlib import Path
import shutil
import time
import re

from . import llm_handler
from .llm.utils import parse_llm_response_for_xml_changes

# Editing and PowerPoint helpers
from .ppt import (
    pptx_to_json,
    extract_xml_from_pptx,
    create_modified_pptx,
    validate_pptx_package,
    validate_xml,
    attempt_repair_xml,
)

# --- Constants ---
MODIFIED_PPTX_FOLDER = Path(
    os.environ.get("VOICE_PPT_OUTPUT_DIR", str(Path.cwd() / "outputs" / "ppt"))
).expanduser()


def _edited_output_path(
    original_pptx_path: str,
    round_number: int = 1,
    output_filepath: str | Path | None = None,
) -> Path:
    if output_filepath:
        path = Path(output_filepath)
        path.parent.mkdir(parents=True, exist_ok=True)
        return path

    stem = Path(original_pptx_path).stem
    stem = re.sub(r'^modified_\d+_', '', stem)
    stem = re.sub(r'_r\d+_edited$', '', stem)
    stem = re.sub(r'_edited$', '', stem)
    MODIFIED_PPTX_FOLDER.mkdir(parents=True, exist_ok=True)
    return MODIFIED_PPTX_FOLDER / f"{stem}_r{round_number}_edited.pptx"

def decide_editing_strategy(
    user_prompt: str,
    ppt_json_data: dict,
    api_key: str,
    request_id: str,
) -> str:
    """
    Uses a preliminary LLM call to decide which editing path to take.

    Args:
        user_prompt: The user's instruction.
        ppt_json_data: The JSON representation of the presentation.
        api_key: The API key to use for the LLM call.
        request_id: The ID for the current request for logging.

    Returns:
        A string, either "XML_EDIT" or "PYTHON_PPTX_EDIT".
    """
    return llm_handler.call_llm_router(
        user_prompt=user_prompt,
        ppt_json_data=ppt_json_data,
        api_key=api_key,
        request_id=request_id,
    )

def _count_slides_edited(original_path: str, modified_path: str) -> int:
    """
    Compare two PPTX files by XML and count how many slides differ.
    Detects text, formatting, and structural changes.
    """
    try:
        from pptx import Presentation
        from lxml import etree
        orig = Presentation(original_path)
        mod = Presentation(modified_path)

        changed = 0
        for i, (orig_slide, mod_slide) in enumerate(zip(orig.slides, mod.slides)):
            orig_xml = etree.tostring(orig_slide._element, encoding="unicode")
            mod_xml = etree.tostring(mod_slide._element, encoding="unicode")
            if orig_xml != mod_xml:
                changed += 1

        # Also count if slide count changed
        slide_diff = abs(len(orig.slides) - len(mod.slides))
        return changed + slide_diff
    except (KeyboardInterrupt, SystemExit):
        raise
    except Exception as e:
        print(f"[Orchestrator] Could not count slides edited: {e}", file=sys.stderr)
        return -1


def _execute_python_pptx_edit(
    original_filepath: str,
    user_prompt: str,
    ppt_json_data: dict,
    api_key: str,
    output_round: int = 1,
    output_filepath: str | Path | None = None,
    trajectory=None,
):
    """
    Manages the two-step LLM chain for python-pptx editing.
    1. Generate content.
    2. Generate code to apply content.
    3. Securely execute the code.
    """
    t_start = time.time()
    t_llm = 0.0

    # Step 1: describe the requested edit as structured content.
    print(f"  {'plan':<8}", end="", flush=True)
    t1 = time.time()
    generated_content = llm_handler.generate_content_for_python_pptx(
        user_prompt=user_prompt,
        ppt_json_data=ppt_json_data,
        api_key=api_key,
        request_id=None,
    )
    t_llm += time.time() - t1

    if not generated_content or "error" in generated_content:
        err = f"Failed to generate content: {generated_content.get('error', 'Unknown error')}"
        if trajectory:
            trajectory.log("execute", "python_pptx", "content_gen_failed",
                           "python-pptx content planning failed — escalating to XML path",
                           error=err, badge_label="PY❌")
        return {"error": err}

    if trajectory:
        trajectory.log("plan", "python_pptx", "content_plan",
                       "python-pptx content plan generated",
                       plan_text=json.dumps(generated_content, indent=2)[:2000],
                       badge_label="PY PLAN")

    # Step 2: generate the python-pptx edit script.
    print(f"  {'code':<8}", end="", flush=True)
    t2 = time.time()
    generated_code = llm_handler.generate_python_pptx_code(
        user_prompt=user_prompt,
        ppt_json_data=ppt_json_data,
        generated_content=generated_content,
        api_key=api_key,
        request_id=None,
    )
    t_llm += time.time() - t2

    if not generated_code or generated_code.strip().startswith("print('Error"):
        err = f"Failed to generate code: {generated_code}"
        if trajectory:
            trajectory.log("execute", "python_pptx", "code_gen_failed",
                           "python-pptx code generation failed — escalating to XML path",
                           error=err, badge_label="PY❌")
        return {"error": err}

    # Step 3: execute the script after the caller's explicit opt-in.
    modified_pptx_path = _execute_generated_code(
        original_filepath,
        generated_code,
        generated_content,
        output_round=output_round,
        output_filepath=output_filepath,
    )

    # Count how many slides were actually edited
    slides_edited = 0
    if modified_pptx_path and Path(modified_pptx_path).exists():
        slides_edited = _count_slides_edited(original_filepath, modified_pptx_path)
    print(f"  exec → {slides_edited} slides")

    if trajectory:
        if modified_pptx_path:
            trajectory.log("execute", "python_pptx", "code_executed",
                           "python-pptx script executed — partial edits applied",
                           before_pptx=original_filepath,
                           after_pptx=modified_pptx_path,
                           badge_label="PY EDIT",
                           extra={"code_snippet": generated_code[:500],
                                  "slides_edited": slides_edited})
        else:
            trajectory.log("execute", "python_pptx", "exec_failed",
                           "python-pptx execution failed — core theme edit needs XML",
                           error="Script execution returned None",
                           badge_label="PY❌")

    result = {
        "modified_pptx_filepath": modified_pptx_path,
        "generated_code": generated_code,
        "generated_content": generated_content,
        "timing_stats": {
            "total_processing_time_s": round(time.time() - t_start, 3),
            "llm_inference_time_s": round(t_llm, 3),
            "number_of_slides_edited_by_llm": slides_edited,
        },
    }
    if not modified_pptx_path:
        result["error"] = "python-pptx script execution failed (no output file produced)"
    return result


def _execute_generated_code(
    original_pptx_path: str,
    code: str,
    content: dict,
    output_round: int = 1,
    output_filepath: str | Path | None = None,
) -> str | None:
    """Execute model-generated Python after the caller explicitly opts in.

    The temporary directory limits file clutter, but this is *not* a security
    sandbox. The child process has the permissions of the current user. The
    public CLI therefore disables this route unless --allow-generated-code is
    supplied.
    """
    final_modified_path = _edited_output_path(original_pptx_path, output_round, output_filepath)

    with tempfile.TemporaryDirectory(prefix="voice_ppt_exec_") as temp_dir_text:
        temp_dir = Path(temp_dir_text)
        same_path = final_modified_path.resolve() == Path(original_pptx_path).resolve()
        backup_path = temp_dir / "input_backup.pptx"
        if same_path:
            shutil.copy2(original_pptx_path, backup_path)
        else:
            shutil.copy2(original_pptx_path, final_modified_path)

        script_path = temp_dir / "generated_edit.py"
        content_path = temp_dir / "content.json"
        script_path.write_text(code, encoding="utf-8")
        content_path.write_text(json.dumps(content, ensure_ascii=False), encoding="utf-8")

        try:
            timeout = int(os.environ.get("VOICE_PPT_CODE_TIMEOUT_S", "120"))
        except ValueError:
            timeout = 120
        timeout = min(max(timeout, 1), 600)

        # This is risk reduction, not a sandbox: isolate Python import state and
        # avoid forwarding API keys or unrelated process secrets to generated
        # code. The child can still access files and the network as the user.
        child_env = {
            key: value
            for key, value in os.environ.items()
            if key in {"PATH", "LANG", "LC_ALL", "LC_CTYPE", "TMPDIR", "TEMP", "TMP"}
        }
        child_env["PYTHONNOUSERSITE"] = "1"
        child_env["PYTHONIOENCODING"] = "utf-8"
        try:
            result = subprocess.run(
                [
                    sys.executable,
                    "-I",
                    str(script_path),
                    str(final_modified_path),
                    str(content_path),
                ],
                cwd=temp_dir,
                env=child_env,
                capture_output=True,
                text=True,
                timeout=timeout,
                check=False,
            )
        except subprocess.TimeoutExpired:
            result = None
            print(f"\n  ✗ generated script timed out after {timeout}s")

        if (
            result is not None
            and result.returncode == 0
            and final_modified_path.exists()
            and validate_pptx_package(final_modified_path)
        ):
            return str(final_modified_path)

        if result is not None and result.returncode == 0 and final_modified_path.exists():
            print("\n  ✗ generated script produced an invalid PPTX package")

        if result is not None and result.returncode != 0:
            err_lines = result.stderr.strip().split("\n")
            print("\n  ✗ script error:")
            for line in err_lines[-15:]:
                print(f"    {line}")

        if same_path and backup_path.exists():
            shutil.copy2(backup_path, final_modified_path)
        elif not same_path and final_modified_path.exists():
            final_modified_path.unlink()
        return None


def _execute_xml_edit(
    original_filepath: str,
    prompt_text: str,
    use_pre_analysis: bool,
    request_id: str,
    api_key: str,
    output_round: int = 1,
    output_filepath: str | Path | None = None,
    ppt_json_data: dict | None = None,
    edit_history=None,
    trajectory=None,
):
    """
    The original XML processing logic, now housed in the orchestrator.
    """
    overall_start_time = time.time()
    original_filename_secure = Path(original_filepath).name
    original_xml_output_dir = None

    try:
        time_json_start = time.time()
        json_data = (
            ppt_json_data
            if isinstance(ppt_json_data, dict)
            else pptx_to_json(original_filepath)
        )
        time_json_end = time.time()

        time_xml_extract_start = time.time()
        original_xml_output_dir = tempfile.mkdtemp(prefix="pptx_xml_")
        extracted_original_xml_full_paths = extract_xml_from_pptx(original_filepath, original_xml_output_dir)
        time_xml_extract_end = time.time()

        print(f"  planning XML...", end="", flush=True)
        if trajectory:
            trajectory.log("plan", "xml", "xml_planning",
                           "LLM analysing PPTX structure and planning XML edits",
                           after_pptx=original_filepath, badge_label="XML PLAN")
        llm_result = llm_handler.get_llm_response(
            user_prompt=prompt_text,
            ppt_json_data=json_data,
            xml_file_paths=extracted_original_xml_full_paths,
            use_pre_analysis=use_pre_analysis,
            request_id=request_id,
            api_key=api_key,
            edit_history=edit_history,
        )
        actual_model_used = llm_result.get("model_used")
        if llm_result.get("error"):
            return {"error": f"LLM request failed: {llm_result['error']}"}
        if trajectory:
            plan_txt = llm_result.get("planning_plan", "") or ""
            trajectory.log("plan", "xml", "xml_plan_received",
                           f"XML edit plan received ({len(plan_txt)} chars)",
                           plan_text=plan_txt[:3000] if plan_txt else None,
                           badge_label="XML PLAN")
        parsed_modified_xml_map = parse_llm_response_for_xml_changes(
            llm_result.get("text_response", "")
        )

        modified_pptx_filepath = None
        number_of_slides_edited = 0
        reason_for_no_modification = None
        time_pptx_modify_start = time_pptx_modify_end = 0

        if parsed_modified_xml_map:
            # Validate and attempt repair for each XML; drop unrecoverable
            repaired_map = {}
            skipped_files = []
            for fname, xml_text in parsed_modified_xml_map.items():
                if validate_xml(xml_text):
                    repaired_map[fname] = xml_text
                    continue
                fixed = attempt_repair_xml(xml_text)
                if fixed and validate_xml(fixed):
                    repaired_map[fname] = fixed
                else:
                    skipped_files.append(fname)

            if skipped_files:
                print(f"\n  ⚠ {len(skipped_files)} invalid XML file(s) skipped: {', '.join(skipped_files)}")
            parsed_modified_xml_map = repaired_map

            if not parsed_modified_xml_map:
                reason_for_no_modification = "All proposed XML changes were invalid and could not be repaired."
                modified_pptx_filepath = None
                # fall through to response payload
            else:
                # Determine which slides were edited
                edited_slide_numbers = set()
                global_change_detected = False
                for llm_filename_key in parsed_modified_xml_map.keys():
                    match = re.search(r'ppt/slides/slide(\d+)\.xml', llm_filename_key)
                    if match:
                        edited_slide_numbers.add(int(match.group(1)))

                    if 'ppt/theme/' in llm_filename_key or 'ppt/slideMasters/' in llm_filename_key or 'ppt/slideLayouts/' in llm_filename_key:
                        global_change_detected = True

                if global_change_detected:
                    total_slides = len(json_data.get("slides", []))
                    for i in range(1, total_slides + 1):
                        edited_slide_numbers.add(i)

                number_of_slides_edited = len(edited_slide_numbers)

                modified_pptx_filepath = _edited_output_path(
                    original_filename_secure,
                    output_round,
                    output_filepath,
                )

                print(f"  applying...", end="", flush=True)
                if trajectory:
                    touched_slides = sorted(list(edited_slide_numbers))
                    changed_xml_files = list(parsed_modified_xml_map.keys())
                    summary = f"Patching {len(changed_xml_files)} XML file(s): {', '.join(changed_xml_files[:3])}"
                    if len(changed_xml_files) > 3:
                        summary += f" … (+{len(changed_xml_files)-3} more)"
                    trajectory.log("execute", "xml", "xml_patch_applied",
                                   summary,
                                   after_pptx=None,
                                   before_pptx=original_filepath,
                                   slides_touched=touched_slides,
                                   badge_label="XML EDIT",
                                   extra={"xml_files_changed": changed_xml_files})
                time_pptx_modify_start = time.time()
                creation_success = create_modified_pptx(
                    original_filepath,
                    parsed_modified_xml_map,
                    str(modified_pptx_filepath)
                )
                time_pptx_modify_end = time.time()

                if creation_success:
                    if trajectory:
                        primary_xml = next(
                            (f for f in parsed_modified_xml_map if "theme" in f or "slide" in f),
                            next(iter(parsed_modified_xml_map), None)
                        )
                        trajectory.log("verify", "xml", "xml_verified",
                                       f"Modified PPTX created — {len(edited_slide_numbers)} slide(s) changed",
                                       before_pptx=original_filepath,
                                       after_pptx=str(modified_pptx_filepath),
                                       slides_touched=sorted(list(edited_slide_numbers)),
                                       xml_diff_file=primary_xml,
                                       verification={"pptx_valid": True, "slides_changed": len(edited_slide_numbers)},
                                       badge_label="VERIFY")
                else:
                    reason_for_no_modification = "PPTX creation failed in ppt_processor."
                    modified_pptx_filepath = None

        else:
            reason_for_no_modification = llm_result.get("text_response", "The LLM did not return any parsable XML modifications.")

        total_processing_time = time.time() - overall_start_time

        timing_stats = {
            "total_processing_time_s": round(total_processing_time, 3),
            "json_extraction_time_s": round(time_json_end - time_json_start, 3),
            "xml_extraction_time_s": round(time_xml_extract_end - time_xml_extract_start, 3),
            "llm_inference_time_s": llm_result.get("inference_time_seconds"),
            "pptx_modification_time_s": round(time_pptx_modify_end - time_pptx_modify_start, 3) if time_pptx_modify_start else "N/A",
            "number_of_slides_edited_by_llm": number_of_slides_edited,
            "total_slides_in_original": len(json_data.get("slides", []))
        }

        response_payload = {
            "message": "File processed successfully (XML Path).",
            "llm_engine_used": actual_model_used,
            "llm_response": llm_result.get("text_response"),
            "reason_for_no_modification": reason_for_no_modification,
            "timing_stats": timing_stats,
            "json_data": json_data,
            "xml_files": [Path(f).name for f in llm_result.get("relevant_files", [])],
            "modified_xml_data": parsed_modified_xml_map,
            "modified_pptx_filepath": str(modified_pptx_filepath) if modified_pptx_filepath else None,
            "planning_plan": llm_result.get("planning_plan"),
            "planning_model": llm_result.get("planning_model"),
        }
        return response_payload

    except Exception as e:
        print(f"Error in _execute_xml_edit: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        return {"error": f"An error occurred during XML processing: {str(e)}"}
    finally:
        if original_xml_output_dir:
            shutil.rmtree(original_xml_output_dir, ignore_errors=True)

# Public entry point used by the CLI.
def _process_single_iteration(
    original_filepath: str,
    prompt_text: str,
    use_pre_analysis: bool = True,
    request_id: str = "",
    api_key: str = None,
    output_round: int = 1,
    output_filepath: str | Path | None = None,
    edit_history=None,
    force_python_pptx: bool = False,
    force_xml: bool = False,
    allow_generated_code: bool = False,
    trajectory=None,
):
    """Run a single pass of the hybrid pipeline (XML vs python-pptx)."""
    ppt_json_data = pptx_to_json(original_filepath)
    if isinstance(ppt_json_data, dict):
        ppt_json_data['request_id'] = request_id

    def _run_strategy(prompt_override: str, chosen_strategy: str):
        if chosen_strategy == "PYTHON_PPTX_EDIT":
            return _execute_python_pptx_edit(
                original_filepath=original_filepath,
                user_prompt=prompt_override,
                ppt_json_data=ppt_json_data,
                api_key=api_key,
                output_round=output_round,
                output_filepath=output_filepath,
                trajectory=trajectory,
            )
        return _execute_xml_edit(
            original_filepath=original_filepath,
            prompt_text=prompt_override,
            use_pre_analysis=use_pre_analysis,
            request_id=request_id,
            api_key=api_key,
            output_round=output_round,
            output_filepath=output_filepath,
            ppt_json_data=ppt_json_data,
            edit_history=edit_history,
            trajectory=trajectory,
        )

    if force_python_pptx and force_xml:
        return {"error": "force_python_pptx and force_xml are mutually exclusive."}

    if force_python_pptx:
        strategy = "PYTHON_PPTX_EDIT"
        if trajectory:
            trajectory.log("route", "router", "force_python_pptx",
                           "Router bypassed — forced python-pptx path",
                           after_pptx=original_filepath,
                           badge_label="ROUTE")
    elif force_xml:
        strategy = "XML_EDIT"
        if trajectory:
            trajectory.log("route", "router", "force_xml",
                           "Router bypassed — forced XML path",
                           after_pptx=original_filepath,
                           badge_label="ROUTE")
    else:
        strategy = decide_editing_strategy(prompt_text, ppt_json_data, api_key, request_id)
        if trajectory:
            trajectory.log("route", "router", "router_decision",
                           f"Router decided: {strategy}",
                           after_pptx=original_filepath,
                           badge_label="ROUTE")
    strategy_label = "python-pptx" if strategy == "PYTHON_PPTX_EDIT" else "XML"
    print(f"  route → {strategy_label}")

    if strategy == "PYTHON_PPTX_EDIT" and not allow_generated_code:
        return {
            "error": (
                "The router selected the generated-Python route, which is disabled by default. "
                "Re-run only in an isolated environment with --allow-generated-code, or use "
                "--execution-mode xml."
            )
        }

    result = _run_strategy(prompt_text, strategy)

    # Retry once on the same pathway, injecting the error context into the prompt
    if isinstance(result, dict) and result.get("error"):
        err_text = str(result.get("error"))
        retry_prompt = f"{prompt_text}\n\n[Retry after error]\nPrevious attempt error: {err_text}\nPlease resolve the issue and try again using the same approach."
        print(f"  ↳ retrying after error...")
        retry_result = _run_strategy(retry_prompt, strategy)
        if isinstance(retry_result, dict):
            retry_result["retry_from_error"] = True
            retry_result["previous_error"] = err_text
        return retry_result

    return result


def process_presentation_hybrid(
    original_filepath: str,
    prompt_text: str,
    use_pre_analysis: bool = True,
    request_id: str = "",
    api_key: str = None,
    edit_history=None,
    force_python_pptx: bool = False,
    force_xml: bool = False,
    allow_generated_code: bool = False,
    loop_mode: bool = False,
    loop_max_iterations: int = 1,
    output_filepath: str | Path | None = None,
    trajectory=None,
):
    """Execute the hybrid pipeline once or looped up to the requested iterations."""

    normalized_iterations = max(1, int(loop_max_iterations or 1))
    use_loop = loop_mode and normalized_iterations > 1

    if not use_loop:
        return _process_single_iteration(
            original_filepath=original_filepath,
            prompt_text=prompt_text,
            use_pre_analysis=use_pre_analysis,
            request_id=request_id,
            api_key=api_key,
            output_round=1,
            output_filepath=output_filepath,
            edit_history=edit_history,
            force_python_pptx=force_python_pptx,
            force_xml=force_xml,
            allow_generated_code=allow_generated_code,
            trajectory=trajectory,
        )

    iteration_summaries = []
    current_input = str(original_filepath)
    final_result = None

    for iteration_idx in range(1, normalized_iterations + 1):
        iter_tag = f"[{iteration_idx}/{normalized_iterations}]" if use_loop else ""
        print(f"\n▶ {iter_tag} Round {iteration_idx}  ".ljust(50, "─"))
        if trajectory:
            trajectory.log("reflect", "xml", f"loop_iter_{iteration_idx}",
                           f"Reflection pass {iteration_idx}/{normalized_iterations}: re-evaluating from {Path(current_input).name}",
                           after_pptx=current_input,
                           badge_label=f"LOOP {iteration_idx}")
        iteration_result = _process_single_iteration(
            original_filepath=current_input,
            prompt_text=prompt_text,
            use_pre_analysis=use_pre_analysis,
            request_id=request_id,
            api_key=api_key,
            output_round=iteration_idx,
            output_filepath=output_filepath if iteration_idx == normalized_iterations else None,
            edit_history=edit_history,
            force_python_pptx=force_python_pptx,
            force_xml=force_xml,
            allow_generated_code=allow_generated_code,
            trajectory=trajectory,
        )

        iter_timing = iteration_result.get("timing_stats", {}) if isinstance(iteration_result, dict) else {}
        summary_entry = {
            "iteration": iteration_idx,
            "input_filepath": current_input,
            "output_filepath": iteration_result.get("modified_pptx_filepath") if isinstance(iteration_result, dict) else None,
            "error": iteration_result.get("error") if isinstance(iteration_result, dict) else "Unknown error",
            "message": iteration_result.get("message") if isinstance(iteration_result, dict) else None,
            "slides_edited": iter_timing.get("number_of_slides_edited_by_llm", "N/A"),
            "elapsed_s": iter_timing.get("total_processing_time_s", "N/A"),
            "llm_s": iter_timing.get("llm_inference_time_s", "N/A"),
        }
        iteration_summaries.append(summary_entry)
        final_result = iteration_result

        if not isinstance(iteration_result, dict):
            print(f"[{request_id}] Loop iteration {iteration_idx} returned non-dict result; aborting loop.")
            break

        if iteration_result.get("error"):
            print(f"[{request_id}] Loop iteration {iteration_idx} failed; stopping further iterations.")
            break

        next_input = iteration_result.get("modified_pptx_filepath")
        if not next_input:
            print(f"[{request_id}] Loop iteration {iteration_idx} produced no PPT output; stopping loop early.")
            break

        current_input = str(next_input)

    if not isinstance(final_result, dict):
        return {"error": "Loop execution did not return a valid result."}

    final_result["loop_mode_enabled"] = True
    final_result["loop_iterations_requested"] = normalized_iterations
    final_result["loop_iterations_completed"] = len(iteration_summaries)
    final_result["loop_iteration_summaries"] = iteration_summaries
    final_result["loop_final_input_filepath"] = current_input
    return final_result
