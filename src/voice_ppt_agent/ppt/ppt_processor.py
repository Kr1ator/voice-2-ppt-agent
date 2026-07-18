# --- ppt_processor.py ---
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pptx.oxml.simpletypes as _pptx_simpletypes
import os

from pathlib import Path, PurePosixPath
import zipfile
import shutil

# ---------------------------------------------------------------------------
# Monkey patch python-pptx integer simple type parsing to tolerate float literals
# (e.g., "6096000.0" inside malformed PPTX files).
# ---------------------------------------------------------------------------
_orig_baseint_convert = _pptx_simpletypes.BaseIntType.convert_from_xml


def _baseint_convert_from_xml(cls, str_value):
    try:
        return _orig_baseint_convert(str_value)
    except ValueError:
        try:
            return int(float(str_value))
        except ValueError:
            raise


_pptx_simpletypes.BaseIntType.convert_from_xml = classmethod(_baseint_convert_from_xml)

_orig_stcoord_convert = _pptx_simpletypes.ST_Coordinate.convert_from_xml


def _stcoord_convert_from_xml(cls, str_value):
    try:
        return _orig_stcoord_convert(str_value)
    except ValueError:
        try:
            if "i" in str_value or "m" in str_value or "p" in str_value:
                return _pptx_simpletypes.ST_UniversalMeasure.convert_from_xml(str_value)
            return _pptx_simpletypes.Emu(int(float(str_value)))
        except ValueError:
            raise


_pptx_simpletypes.ST_Coordinate.convert_from_xml = classmethod(_stcoord_convert_from_xml)


def _normalize_pptx_xml_path(path: str) -> str:
    """Return the PPTX-internal XML path when an extracted temp path is provided."""
    normalized = str(path or "").strip().strip('"').strip("'").replace("\\", "/")
    if normalized.startswith("file://"):
        normalized = normalized[len("file://"):]

    for marker in ("ppt/", "docProps/", "_rels/"):
        idx = normalized.find(marker)
        if idx != -1:
            return normalized[idx:]

    content_types = "[Content_Types].xml"
    idx = normalized.find(content_types)
    if idx != -1:
        return content_types

    return normalized


def _validate_internal_xml_path(path: str) -> str:
    """Validate a normalized path before it is written into a PPTX archive."""
    candidate = PurePosixPath(path)
    if candidate.is_absolute() or ".." in candidate.parts or not candidate.parts:
        raise ValueError(f"Unsafe modified PPTX path: {path}")
    if path != "[Content_Types].xml" and not path.endswith((".xml", ".rels")):
        raise ValueError(f"Only PPTX XML and relationship parts may be modified: {path}")
    return candidate.as_posix()

def extract_text_from_shape(shape):
    """Extracts text from a shape, handling different shape types."""
    text = ""
    if shape.has_text_frame:
        text = shape.text_frame.text
    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text += cell.text_frame.text + "\t"
            text += "\n"
    return text.strip()

def pptx_to_json(filepath):
    """
    Converts a .pptx file to a comprehensive JSON representation.
    Captures all edit types: Content, Layout, Styling, Interactivity, and Structure.
    """
    try:
        prs = Presentation(filepath)

        # Presentation-level metadata
        presentation_data = {
            "filename": os.path.basename(filepath),
            "slide_width": prs.slide_width.pt if hasattr(prs, 'slide_width') else None,
            "slide_height": prs.slide_height.pt if hasattr(prs, 'slide_height') else None,
            "slides": []
        }

        # LOW PRIORITY: Custom properties/metadata (existence check)
        try:
            if hasattr(prs, 'core_properties'):
                props = prs.core_properties
                custom_props = {}
                if hasattr(props, 'author') and props.author:
                    custom_props["author"] = props.author
                if hasattr(props, 'title') and props.title:
                    custom_props["title"] = props.title
                if hasattr(props, 'subject') and props.subject:
                    custom_props["subject"] = props.subject
                if hasattr(props, 'keywords') and props.keywords:
                    custom_props["keywords"] = props.keywords
                if custom_props:
                    presentation_data["custom_properties"] = custom_props
        except Exception:            pass

        for i, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": i + 1,
                "shapes": [],
                "notes": "",
                # Layout & Structure info
                "slide_layout": slide.slide_layout.name if hasattr(slide, 'slide_layout') else None,
                "slide_id": slide.slide_id if hasattr(slide, 'slide_id') else None,
            }

            # Background info (Styling category)
            try:
                if hasattr(slide, 'background'):
                    bg = slide.background
                    slide_data["background"] = {
                        "fill_type": str(bg.fill.type) if hasattr(bg, 'fill') else None,
                    }
                    try:
                        if hasattr(bg.fill, 'fore_color') and bg.fill.fore_color.rgb:
                            slide_data["background"]["color_rgb"] = str(bg.fill.fore_color.rgb)
                    except Exception:                        pass
            except Exception:
                pass

            for shape_idx, shape in enumerate(slide.shapes):
                shape_info = {
                    "shape_id": shape.shape_id if hasattr(shape, 'shape_id') else None,
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "text": extract_text_from_shape(shape),
                    # Layout info (position, size, rotation)
                    "left": shape.left.pt if hasattr(shape, 'left') and shape.left is not None else None,
                    "top": shape.top.pt if hasattr(shape, 'top') and shape.top is not None else None,
                    "width": shape.width.pt if hasattr(shape, 'width') and shape.width is not None else None,
                    "height": shape.height.pt if hasattr(shape, 'height') and shape.height is not None else None,
                    "rotation": shape.rotation if hasattr(shape, 'rotation') else None,
                    "z_order": shape_idx,  # Z-order based on iteration order
                }

                # Placeholder info (Layout category)
                try:
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        shape_info["is_placeholder"] = True
                        shape_info["placeholder_type"] = str(shape.placeholder_format.type) if hasattr(shape, 'placeholder_format') else None
                except Exception:                    pass

                # Shape styling (fills, lines, effects)
                try:
                    if hasattr(shape, 'fill'):
                        fill_info = {
                            "type": str(shape.fill.type) if shape.fill else None,
                        }
                        try:
                            if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color and shape.fill.fore_color.rgb:
                                fill_info["color_rgb"] = str(shape.fill.fore_color.rgb)
                        except Exception:                            pass

                        # HIGH PRIORITY: Gradient fills
                        try:
                            if hasattr(shape.fill, 'gradient_angle'):
                                fill_info["gradient_angle"] = shape.fill.gradient_angle
                            if hasattr(shape.fill, 'gradient_stops'):
                                stops = []
                                for stop in shape.fill.gradient_stops:
                                    stop_info = {
                                        "position": stop.position if hasattr(stop, 'position') else None,
                                    }
                                    try:
                                        if hasattr(stop, 'color') and stop.color and stop.color.rgb:
                                            stop_info["color_rgb"] = str(stop.color.rgb)
                                    except Exception:                                        pass
                                    stops.append(stop_info)
                                if stops:
                                    fill_info["gradient_stops"] = stops
                        except Exception:                            pass

                        shape_info["fill"] = fill_info
                except Exception:                    pass

                try:
                    if hasattr(shape, 'line'):
                        line_info = {
                            "width": shape.line.width.pt if shape.line.width else None,
                        }
                        try:
                            if hasattr(shape.line, 'color') and shape.line.color and shape.line.color.rgb:
                                line_info["color_rgb"] = str(shape.line.color.rgb)
                        except Exception:                            pass
                        # HIGH PRIORITY: Line dash style
                        try:
                            if hasattr(shape.line, 'dash_style') and shape.line.dash_style:
                                line_info["dash_style"] = str(shape.line.dash_style)
                        except Exception:                            pass
                        shape_info["line"] = line_info
                except Exception:                    pass

                # Shadow effect
                try:
                    if hasattr(shape, 'shadow') and shape.shadow and shape.shadow.inherit:
                        shape_info["has_shadow"] = True
                except Exception:                    pass

                # Hyperlinks (Interactivity category)
                try:
                    if hasattr(shape, 'click_action') and shape.click_action and shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                        shape_info["hyperlink"] = shape.click_action.hyperlink.address
                except Exception:                    pass

                # Extract detailed TYPOGRAPHY from text frames (Content category)
                if shape.has_text_frame:
                    shape_info["paragraphs"] = []
                    for para in shape.text_frame.paragraphs:
                        para_info = {
                            "text": para.text,
                            "level": para.level,
                            "alignment": str(para.alignment) if hasattr(para, 'alignment') else None,
                            "line_spacing": para.line_spacing if hasattr(para, 'line_spacing') else None,
                            "space_before": para.space_before.pt if hasattr(para, 'space_before') and para.space_before else None,
                            "space_after": para.space_after.pt if hasattr(para, 'space_after') and para.space_after else None,
                            "runs": []
                        }

                        # HIGH PRIORITY: Bullet/Numbering formats
                        try:
                            if hasattr(para, 'bullet_format') and para.bullet_format:
                                bullet_info = {}
                                if hasattr(para.bullet_format, 'type') and para.bullet_format.type:
                                    bullet_info["type"] = str(para.bullet_format.type)
                                if hasattr(para.bullet_format, 'char') and para.bullet_format.char:
                                    bullet_info["char"] = para.bullet_format.char
                                if hasattr(para.bullet_format, 'start_value') and para.bullet_format.start_value:
                                    bullet_info["start_value"] = para.bullet_format.start_value
                                if bullet_info:
                                    para_info["bullet_format"] = bullet_info
                        except Exception:                            pass

                        for run in para.runs:
                            run_info = {
                                "text": run.text,
                                "font_name": run.font.name if run.font.name else None,
                                "font_size": run.font.size.pt if run.font.size else None,
                                "bold": run.font.bold,
                                "italic": run.font.italic,
                                "underline": run.font.underline,
                            }

                            # Color info
                            try:
                                if run.font.color and run.font.color.rgb:
                                    run_info["color_rgb"] = str(run.font.color.rgb)
                            except Exception:
                                pass

                            # Hyperlink in run
                            try:
                                if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                                    run_info["hyperlink"] = run.hyperlink.address
                            except Exception:                                pass

                            para_info["runs"].append(run_info)

                        shape_info["paragraphs"].append(para_info)

                # TABLES (Content category)
                if shape.has_table:
                    table_info = {
                        "rows": len(shape.table.rows),
                        "columns": len(shape.table.columns),
                        "cells": []
                    }

                    for row_idx, row in enumerate(shape.table.rows):
                        table_info["row_height_" + str(row_idx)] = row.height.pt if row.height else None
                        for col_idx, cell in enumerate(row.cells):
                            cell_info = {
                                "row": row_idx,
                                "col": col_idx,
                                "text": cell.text_frame.text if hasattr(cell, 'text_frame') else "",
                                "paragraphs": []
                            }

                            # Cell fill
                            try:
                                if hasattr(cell, 'fill') and cell.fill:
                                    cell_info["fill_type"] = str(cell.fill.type)
                                    if hasattr(cell.fill, 'fore_color') and cell.fill.fore_color and cell.fill.fore_color.rgb:
                                        cell_info["fill_color_rgb"] = str(cell.fill.fore_color.rgb)
                            except Exception:                                pass

                            # MEDIUM PRIORITY: Table cell borders
                            try:
                                if hasattr(cell, 'borders'):
                                    borders_info = {}
                                    for side in ['top', 'bottom', 'left', 'right']:
                                        if hasattr(cell.borders, side):
                                            border = getattr(cell.borders, side)
                                            if border:
                                                border_info = {}
                                                try:
                                                    if hasattr(border, 'width') and border.width:
                                                        border_info["width"] = border.width.pt
                                                except Exception:                                                    pass
                                                try:
                                                    if hasattr(border, 'color') and border.color and border.color.rgb:
                                                        border_info["color_rgb"] = str(border.color.rgb)
                                                except Exception:                                                    pass
                                                try:
                                                    if hasattr(border, 'dash_style'):
                                                        border_info["dash_style"] = str(border.dash_style)
                                                except Exception:                                                    pass
                                                if border_info:
                                                    borders_info[side] = border_info
                                    if borders_info:
                                        cell_info["borders"] = borders_info
                            except Exception:                                pass

                            # Cell typography
                            if hasattr(cell, 'text_frame'):
                                for para in cell.text_frame.paragraphs:
                                    para_info = {
                                        "text": para.text,
                                        "runs": []
                                    }

                                    for run in para.runs:
                                        run_info = {
                                            "text": run.text,
                                            "font_name": run.font.name if run.font.name else None,
                                            "font_size": run.font.size.pt if run.font.size else None,
                                            "bold": run.font.bold,
                                            "italic": run.font.italic,
                                        }
                                        try:
                                            if run.font.color and run.font.color.rgb:
                                                run_info["color_rgb"] = str(run.font.color.rgb)
                                        except Exception:                                            pass
                                        para_info["runs"].append(run_info)

                                    cell_info["paragraphs"].append(para_info)

                            table_info["cells"].append(cell_info)

                    # Column widths
                    for col_idx, col in enumerate(shape.table.columns):
                        table_info["column_width_" + str(col_idx)] = col.width.pt if col.width else None

                    shape_info["table"] = table_info

                # IMAGES & PICTURES (Content category)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    shape_info["is_picture"] = True
                    try:
                        if hasattr(shape, 'image'):
                            shape_info["image_type"] = shape.image.content_type if hasattr(shape.image, 'content_type') else None
                            shape_info["image_filename"] = shape.image.filename if hasattr(shape.image, 'filename') else None
                    except Exception:                        pass

                # CHARTS (Content category)
                if hasattr(shape, 'has_chart') and shape.has_chart:
                    shape_info["is_chart"] = True
                    try:
                        chart = shape.chart
                        chart_info = {
                            "chart_type": str(chart.chart_type) if hasattr(chart, 'chart_type') else None,
                            "has_title": chart.has_title if hasattr(chart, 'has_title') else None,
                            "chart_title": chart.chart_title.text_frame.text if hasattr(chart, 'chart_title') and chart.has_title else None,
                        }

                        # HIGH PRIORITY: Chart series data
                        try:
                            if hasattr(chart, 'series'):
                                series_list = []
                                for series in chart.series:
                                    series_info = {
                                        "name": series.name if hasattr(series, 'name') else None,
                                    }
                                    # Extract data values if available
                                    try:
                                        if hasattr(series, 'values') and series.values:
                                            series_info["values"] = list(series.values)
                                    except Exception:                                        pass
                                    series_list.append(series_info)
                                if series_list:
                                    chart_info["series"] = series_list
                        except Exception:                            pass

                        # Extract categories if available
                        try:
                            if hasattr(chart, 'categories') and chart.categories:
                                chart_info["categories"] = list(chart.categories)
                        except Exception:                            pass

                        shape_info["chart"] = chart_info
                    except Exception:                        pass

                # GROUP SHAPES (Layout category)
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    shape_info["is_group"] = True
                    shape_info["group_shape_count"] = len(shape.shapes) if hasattr(shape, 'shapes') else 0

                # HIGH PRIORITY: Shape connectors (for flowcharts, org charts)
                try:
                    if hasattr(shape, 'connector_format'):
                        connector_info = {}
                        if hasattr(shape.connector_format, 'begin_connected') and shape.connector_format.begin_connected:
                            connector_info["begin_connected"] = True
                            if hasattr(shape.connector_format, 'begin_connection_site'):
                                connector_info["begin_connection_site"] = shape.connector_format.begin_connection_site
                        if hasattr(shape.connector_format, 'end_connected') and shape.connector_format.end_connected:
                            connector_info["end_connected"] = True
                            if hasattr(shape.connector_format, 'end_connection_site'):
                                connector_info["end_connection_site"] = shape.connector_format.end_connection_site
                        if connector_info:
                            shape_info["connector"] = connector_info
                except Exception:                    pass

                # Alt Text / Accessibility (Structure category)
                try:
                    if hasattr(shape, 'name'):
                        shape_info["alt_text"] = shape.name
                except Exception:                    pass

                slide_data["shapes"].append(shape_info)

            # NOTES (Structure category)
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                slide_data["notes"] = text_frame.text.strip()

            # SLIDE TRANSITIONS (Interactivity category) - LOW PRIORITY: Enhanced detection
            try:
                # Check for transition via element property
                if hasattr(slide, 'element') and hasattr(slide.element, 'transition'):
                    transition = slide.element.transition
                    if transition is not None:
                        slide_data["has_transition"] = True
                        # Try to extract transition type if available
                        try:
                            if hasattr(transition, 'type'):
                                slide_data["transition_type"] = str(transition.type)
                        except Exception:                            pass
                elif hasattr(slide, 'slide_transitions'):
                    slide_data["has_transition"] = True
            except Exception:                pass

            # LOW PRIORITY: SmartArt detection (existence check)
            try:
                has_smartart = False
                for shape in slide.shapes:
                    # SmartArt is typically detected as a group with diagram data
                    if hasattr(shape, 'element') and hasattr(shape.element, 'tag'):
                        tag = str(shape.element.tag)
                        if 'dgm' in tag.lower() or 'diagram' in tag.lower():
                            has_smartart = True
                            break
                if has_smartart:
                    slide_data["has_smartart"] = True
            except Exception:                pass

            presentation_data["slides"].append(slide_data)

        return presentation_data
    except Exception as e:
        print(f"Error converting {filepath} to JSON: {e}")
        raise



def extract_xml_from_pptx(pptx_filepath, output_folder):
    """
    Extracts all constituent XML files from a .pptx file.
    Returns a list of full paths to the extracted XML files.
    """
    extracted_files_paths = []
    try:
        output_root = Path(output_folder).resolve()
        output_root.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(pptx_filepath, 'r') as pptx_zip:
            for member_info in pptx_zip.infolist():
                member_name = member_info.filename
                if not member_info.is_dir() and member_name.endswith(('.xml', '.rels')):
                    member_path = PurePosixPath(member_name)
                    if member_path.is_absolute() or ".." in member_path.parts:
                        raise ValueError(f"Unsafe path in PPTX archive: {member_name}")
                    target_path = output_root.joinpath(*member_path.parts)
                    target_path.parent.mkdir(parents=True, exist_ok=True)
                    with pptx_zip.open(member_name) as source, target_path.open("wb") as target:
                        shutil.copyfileobj(source, target)
                    extracted_files_paths.append(str(target_path))
        return extracted_files_paths
    except Exception as e:
        print(f"Error extracting XML from {pptx_filepath}: {e}")
        raise

def create_modified_pptx(original_pptx_path, modified_xml_map, output_pptx_path):
    """
    Creates a new .pptx file by taking an original .pptx, and either updating
    existing internal XML files or adding new ones based on the modified_xml_map.
    """
    temp_output_pptx_path = output_pptx_path + ".tmp"
    try:
        Path(output_pptx_path).parent.mkdir(parents=True, exist_ok=True)
        modified_xml_map = {
            _validate_internal_xml_path(_normalize_pptx_xml_path(filename)): content
            for filename, content in modified_xml_map.items()
        }

        # Keep track of which files from the modification map we have used.
        processed_map_files = set()

        with zipfile.ZipFile(original_pptx_path, 'r') as zin:
            with zipfile.ZipFile(temp_output_pptx_path, 'w', zipfile.ZIP_DEFLATED) as zout:
                # 1. Iterate through existing files in the original PPTX.
                for item in zin.infolist():
                    item_name_normalized = item.filename.replace("\\", "/")
                    if item_name_normalized in modified_xml_map:
                        # If an existing file is in our map, write the modified content.
                        new_content = modified_xml_map[item_name_normalized]
                        zout.writestr(item, new_content.encode('utf-8'))
                        processed_map_files.add(item_name_normalized)
                    else:
                        # Otherwise, copy the original file as-is.
                        buffer = zin.read(item.filename)
                        zout.writestr(item, buffer)

                # 2. Iterate through the map to find any *new* files to add.
                for filename, content in modified_xml_map.items():
                    if filename not in processed_map_files:
                        print(f"Adding new file to PPTX archive: {filename}")
                        # writestr can take a ZipInfo object or a string filename.
                        zout.writestr(filename, content.encode('utf-8'))

        os.replace(temp_output_pptx_path, output_pptx_path)
        if not validate_pptx_package(output_pptx_path):
            raise ValueError("The rebuilt file is not a valid PPTX package.")
        print(f"Modified PPTX successfully created at: {output_pptx_path}")
        return True
    except Exception as e:
        print(f"Error creating modified PPTX at {output_pptx_path}: {e}")
        if os.path.exists(temp_output_pptx_path):
            os.remove(temp_output_pptx_path)
        return False


def validate_pptx_package(path: str | Path) -> bool:
    """Perform a fast package-level integrity check without rendering."""
    required = {"[Content_Types].xml", "ppt/presentation.xml"}
    try:
        with zipfile.ZipFile(path, "r") as archive:
            names = set(archive.namelist())
            return required.issubset(names) and archive.testzip() is None
    except (OSError, zipfile.BadZipFile):
        return False



def validate_xml(xml_text):
    """Return True if xml_text is well-formed."""
    try:
        from lxml import etree
        etree.fromstring(xml_text.encode("utf-8"))
        return True
    except Exception as e:
        print(f"Invalid XML detected: {e}")
        return False

def attempt_repair_xml(xml_text: str) -> str | None:
    """
    Attempt to repair malformed XML using lxml's recovery parser.
    Returns a repaired XML string if successful, otherwise None.
    """
    try:
        from lxml import etree
        parser = etree.XMLParser(recover=True)
        root = etree.fromstring(xml_text.encode("utf-8"), parser=parser)
        if root is None:
            return None
        # Serialize back to UTF-8 string
        fixed = etree.tostring(root, encoding="utf-8", xml_declaration=True).decode("utf-8")
        return fixed
    except Exception as e:
        print(f"XML repair failed: {e}")
        return None
