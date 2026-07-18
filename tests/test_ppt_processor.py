from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from pptx import Presentation

from voice_ppt_agent.ppt import (
    create_modified_pptx,
    extract_xml_from_pptx,
    pptx_to_json,
    validate_pptx_package,
)
from voice_ppt_agent import orchestrator


def _make_deck(path: Path, text: str = "Hello SURF") -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = text
    slide.placeholders[1].text = "Offline test fixture"
    presentation.save(path)


def test_inspect_and_rebuild_pptx(tmp_path: Path) -> None:
    original = tmp_path / "input.pptx"
    output = tmp_path / "output.pptx"
    extracted = tmp_path / "xml"
    _make_deck(original)

    summary = pptx_to_json(str(original))
    assert len(summary["slides"]) == 1
    assert any(shape["text"] == "Hello SURF" for shape in summary["slides"][0]["shapes"])
    assert validate_pptx_package(original)

    paths = extract_xml_from_pptx(original, extracted)
    slide_xml_path = next(Path(path) for path in paths if path.endswith("ppt/slides/slide1.xml"))
    changed_xml = slide_xml_path.read_text(encoding="utf-8").replace("Hello SURF", "Goodbye SURF")
    assert create_modified_pptx(
        str(original),
        {str(slide_xml_path): changed_xml},
        str(output),
    )
    assert validate_pptx_package(output)
    assert Presentation(output).slides[0].shapes.title.text == "Goodbye SURF"


def test_extract_rejects_zip_traversal(tmp_path: Path) -> None:
    malicious = tmp_path / "malicious.pptx"
    extraction_root = tmp_path / "extract"
    escaped = tmp_path / "escaped.xml"
    with ZipFile(malicious, "w", ZIP_DEFLATED) as archive:
        archive.writestr("../escaped.xml", "<unsafe />")

    with pytest.raises(ValueError, match="Unsafe path"):
        extract_xml_from_pptx(malicious, extraction_root)
    assert not escaped.exists()


def test_invalid_file_fails_package_check(tmp_path: Path) -> None:
    invalid = tmp_path / "not-a-pptx.pptx"
    invalid.write_bytes(b"not a zip archive")
    assert validate_pptx_package(invalid) is False


def test_rebuild_rejects_unsafe_internal_path(tmp_path: Path) -> None:
    original = tmp_path / "input.pptx"
    output = tmp_path / "output.pptx"
    _make_deck(original)
    assert create_modified_pptx(
        str(original),
        {"../../outside.xml": "<unsafe />"},
        str(output),
    ) is False
    assert not output.exists()


def test_generated_code_output_must_be_a_valid_pptx(tmp_path: Path) -> None:
    original = tmp_path / "input.pptx"
    output = tmp_path / "output.pptx"
    _make_deck(original)
    code = (
        "from pathlib import Path\n"
        "import sys\n"
        "Path(sys.argv[1]).write_bytes(b'not a pptx')\n"
    )

    assert (
        orchestrator._execute_generated_code(
            str(original), code, {}, output_filepath=output
        )
        is None
    )
    assert not output.exists()


def test_generated_code_can_return_a_valid_pptx(tmp_path: Path) -> None:
    original = tmp_path / "input.pptx"
    output = tmp_path / "output.pptx"
    _make_deck(original)
    code = (
        "import sys\n"
        "from pptx import Presentation\n"
        "deck = Presentation(sys.argv[1])\n"
        "deck.slides[0].shapes.title.text = 'Isolated edit'\n"
        "deck.save(sys.argv[1])\n"
    )

    result = orchestrator._execute_generated_code(
        str(original), code, {}, output_filepath=output
    )
    assert result == str(output)
    assert validate_pptx_package(output)
    assert Presentation(output).slides[0].shapes.title.text == "Isolated edit"
